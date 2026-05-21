"""生成《多智能体协作原理》教学PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 配色方案 ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2 = RGBColor(0x7B, 0x68, 0xEE)
ACCENT3 = RGBColor(0x00, 0xE6, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xBB, 0xCC)
ORANGE = RGBColor(0xFF, 0xA5, 0x02)
RED = RGBColor(0xFF, 0x4D, 0x6A)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # 调整圆角
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_text(text_frame, text, font_size=16, color=WHITE, bold=False, level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_after = Pt(4)
    return p


def add_arrow(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ════════════════════════════════════════════
# 第1页：封面
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# 装饰线条
add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Pt(3), fill_color=ACCENT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "多智能体协作原理", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
             "Multi-Agent Collaboration", font_size=28, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
             "以钻井事故处置系统为例，理解智能体如何协同工作", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
             "基于 LangGraph 框架  |  项目实战讲解", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 第2页：什么是多智能体系统
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "什么是多智能体系统？", font_size=36, color=WHITE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.3), Pt(80), Pt(4), fill_color=ACCENT)

# 左侧：定义
card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5), border_color=ACCENT)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(4.8), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "核心定义"
p.font.size = Pt(24)
p.font.color.rgb = ACCENT
p.font.bold = True

bullets = [
    "多个 AI 智能体（Agent）协同完成复杂任务",
    "每个智能体有明确的角色和职责",
    "通过共享状态（State）传递信息",
    "由编排引擎（Orchestrator）协调执行顺序",
    "可以串行、并行或条件分支执行",
]
for b in bullets:
    add_bullet_text(tf, f"  •  {b}", font_size=17, color=WHITE)

# 右侧：类比
card2 = add_shape(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(5), border_color=ACCENT2)
txBox2 = slide.shapes.add_textbox(Inches(7.4), Inches(2.0), Inches(4.8), Inches(4.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "类比：医院会诊"
p2.font.size = Pt(24)
p2.font.color.rgb = ACCENT2
p2.font.bold = True

analogy = [
    "急诊科医生 → 采集信息（事故解析）",
    "影像科 → 查找历史病例（案例匹配）",
    "外科医生 A → 激进方案（激进计划）",
    "外科医生 B → 保守方案（保守计划）",
    "质控部门 → 合规审查（合规检查）",
    "专家组 → 综合决策（最终方案）",
]
for a in analogy:
    add_bullet_text(tf2, f"  •  {a}", font_size=17, color=WHITE)

# ════════════════════════════════════════════
# 第3页：本项目架构总览
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "项目架构：钻井事故处置系统", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Pt(80), Pt(4), fill_color=ACCENT)

# 流程图 - 用卡片和箭头表示
# Row 1: 输入
def make_agent_card(slide, left, top, width, height, title, subtitle, color, border_color=None):
    card = add_shape(slide, left, top, width, height, border_color=border_color or color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.4),
                 title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.45), width - Inches(0.2), height - Inches(0.5),
                 subtitle, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# 第一行：输入 + 事故解析 + 案例匹配
make_agent_card(slide, Inches(0.5), Inches(1.8), Inches(2.2), Inches(1.2),
                "用户输入", "事故文字描述", GRAY, ACCENT)
add_arrow(slide, Inches(2.8), Inches(2.2), Inches(0.6), Inches(0.35), ACCENT)

make_agent_card(slide, Inches(3.5), Inches(1.8), Inches(2.2), Inches(1.2),
                "事故解析器", "Accident Parser\n提取结构化信息", ACCENT, ACCENT)
add_arrow(slide, Inches(5.8), Inches(2.2), Inches(0.6), Inches(0.35), ACCENT)

make_agent_card(slide, Inches(6.5), Inches(1.8), Inches(2.2), Inches(1.2),
                "案例匹配器", "Case Matcher\n检索历史案例", ACCENT, ACCENT)

# 箭头向下分叉
add_down_arrow(slide, Inches(7.3), Inches(3.1), Inches(0.35), Inches(0.5), ACCENT)

# 第二行：并行 - 激进方案 & 保守方案
make_agent_card(slide, Inches(4.2), Inches(3.8), Inches(2.5), Inches(1.2),
                "激进方案 Agent", "Aggressive Plan\n快速恢复、最小停工", ORANGE, ORANGE)

make_agent_card(slide, Inches(8), Inches(3.8), Inches(2.5), Inches(1.2),
                "保守方案 Agent", "Conservative Plan\n安全优先、分阶段处置", ACCENT3, ACCENT3)

# 横向箭头指向合规检查
add_arrow(slide, Inches(6.8), Inches(4.2), Inches(1.1), Inches(0.35), ACCENT)

# 合规检查
make_agent_card(slide, Inches(10.2), Inches(3.8), Inches(2.5), Inches(1.2),
                "合规检查器", "Compliance Checker\n行业标准审查", RED, RED)

# 箭头向下
add_down_arrow(slide, Inches(11.2), Inches(5.1), Inches(0.35), Inches(0.5), ACCENT)

# 第三行：决策 + 输出
make_agent_card(slide, Inches(8.5), Inches(5.8), Inches(2.5), Inches(1.2),
                "决策 maker", "Decision Maker\n综合生成最终方案", ACCENT2, ACCENT2)

add_arrow(slide, Inches(11.1), Inches(6.2), Inches(0.6), Inches(0.35), ACCENT)

make_agent_card(slide, Inches(11.8), Inches(5.8), Inches(1.2), Inches(1.2),
                "输出", "存档", GRAY, ACCENT)

# ════════════════════════════════════════════
# 第4页：LangGraph 框架核心概念
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "LangGraph 核心概念", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Pt(80), Pt(4), fill_color=ACCENT)

concepts = [
    ("StateGraph\n状态图", "定义智能体的执行拓扑\n节点（Node）= 智能体\n边（Edge）= 执行顺序", ACCENT),
    ("State\n共享状态", "所有节点共享的数据结构\nTypedDict 定义字段\n支持 Reducer 合并策略", ACCENT2),
    ("并行执行\nFan-out / Fan-in", "多个节点同时执行\n结果自动合并到 State\n本项目：激进 & 保守并行", ORANGE),
    ("条件路由\nConditional Edge", "根据 State 决定下一步\nrouter 节点判断意图\n分支到不同处理路径", ACCENT3),
]

for i, (title, desc, color) in enumerate(concepts):
    left = Inches(0.8 + i * 3.1)
    card = add_shape(slide, left, Inches(1.8), Inches(2.8), Inches(4.5), border_color=color)
    add_text_box(slide, left + Inches(0.2), Inches(2.0), Inches(2.4), Inches(1.0),
                 title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(3.2), Inches(2.4), Inches(2.8),
                 desc, font_size=15, color=WHITE, alignment=PP_ALIGN.LEFT)

# ════════════════════════════════════════════
# 第5页：State 状态流转
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "共享状态（State）如何流转", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Pt(80), Pt(4), fill_color=ACCENT)

# 左侧：状态定义
card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.2), border_color=ACCENT)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(5.0), Inches(4.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AgentState 关键字段"
p.font.size = Pt(22)
p.font.color.rgb = ACCENT
p.font.bold = True

fields = [
    ("accident", "结构化事故信息（井型、深度、鱼顶类型...）"),
    ("similar_cases", "历史相似案例 Top 3"),
    ("aggressive_plan", "激进处置方案文本"),
    ("conservative_plan", "保守处置方案文本"),
    ("compliance_report", "合规审查报告"),
    ("final_plan", "最终综合方案（Markdown）"),
    ("evidence", "证据链（Annotated + add 累加器）"),
    ("confidence_score", "置信度评分 0.35~0.72"),
]
for name, desc in fields:
    add_bullet_text(tf, f"  {name}", font_size=15, color=ACCENT2, bold=True)
    add_bullet_text(tf, f"    {desc}", font_size=13, color=GRAY, level=1)

# 右侧：流转示意
card2 = add_shape(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(5.2), border_color=ACCENT2)
txBox2 = slide.shapes.add_textbox(Inches(7.5), Inches(2.0), Inches(5.0), Inches(4.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "数据流转过程"
p2.font.size = Pt(22)
p2.font.color.rgb = ACCENT2
p2.font.bold = True

flow = [
    "1. 用户输入 → raw_description",
    "2. 事故解析器 → accident + parse_report",
    "3. 案例匹配 → similar_cases + evidence",
    "4. 激进方案 → aggressive_plan + evidence",
    "5. 保守方案 → conservative_plan + evidence",
    "6. 合规检查 → compliance_report + confidence_score",
    "7. 决策 maker → final_plan",
    "8. 归档 → output_path",
]
for f in flow:
    add_bullet_text(tf2, f"  {f}", font_size=15, color=WHITE)

add_bullet_text(tf2, "", font_size=8, color=WHITE)
add_bullet_text(tf2, "  evidence 字段使用累加器（Reducer）", font_size=14, color=ORANGE, bold=True)
add_bullet_text(tf2, "  并行节点各自追加，自动合并", font_size=14, color=ORANGE)

# ════════════════════════════════════════════
# 第6页：并行执行详解
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "并行执行：两个方案同时生成", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Pt(80), Pt(4), fill_color=ACCENT)

# 中间：案例匹配
make_agent_card(slide, Inches(5.2), Inches(1.6), Inches(2.8), Inches(0.9),
                "案例匹配器", "输出 similar_cases", ACCENT, ACCENT)

# 下方分叉箭头
add_shape(slide, Inches(6.3), Inches(2.55), Pt(3), Inches(0.5), fill_color=ACCENT)
# 左箭头
shape_l = slide.shapes.add_shape(MSO_SHAPE.BENT_ARROW, Inches(3.8), Inches(2.8), Inches(2.5), Inches(0.5))
shape_l.fill.solid()
shape_l.fill.fore_color.rgb = ACCENT
shape_l.line.fill.background()
# 右箭头
shape_r = slide.shapes.add_shape(MSO_SHAPE.BENT_ARROW, Inches(6.8), Inches(2.8), Inches(2.5), Inches(0.5))
shape_r.fill.solid()
shape_r.fill.fore_color.rgb = ACCENT
shape_r.line.fill.background()
shape_r.rotation = 180.0

# 左：激进方案
card_l = add_shape(slide, Inches(1), Inches(3.5), Inches(4.5), Inches(3.2), border_color=ORANGE)
txBox_l = slide.shapes.add_textbox(Inches(1.3), Inches(3.7), Inches(4.0), Inches(2.8))
tf_l = txBox_l.text_frame
tf_l.word_wrap = True
p_l = tf_l.paragraphs[0]
p_l.text = "激进方案 Agent"
p_l.font.size = Pt(22)
p_l.font.color.rgb = ORANGE
p_l.font.bold = True

agg_points = [
    "目标：快速恢复生产，最小化停工时间",
    "策略：直接打捞 → 上击器 → 套铣",
    "特点：步骤少、速度快",
    "风险：可能遗漏复杂情况",
]
for a in agg_points:
    add_bullet_text(tf_l, f"  •  {a}", font_size=15, color=WHITE)

# 右：保守方案
card_r = add_shape(slide, Inches(7.5), Inches(3.5), Inches(4.5), Inches(3.2), border_color=ACCENT3)
txBox_r = slide.shapes.add_textbox(Inches(7.8), Inches(3.7), Inches(4.0), Inches(2.8))
tf_r = txBox_r.text_frame
tf_r.word_wrap = True
p_r = tf_r.paragraphs[0]
p_r.text = "保守方案 Agent"
p_r.font.size = Pt(22)
p_r.font.color.rgb = ACCENT3
p_r.font.bold = True

con_points = [
    "目标：安全优先，防止事故恶化",
    "策略：循环洗井 → 浸泡 → 上击 → 打捞 → 套铣 → 侧钻",
    "特点：分阶段、有停工判定条件",
    "风险：周期长、成本高",
]
for c in con_points:
    add_bullet_text(tf_r, f"  •  {c}", font_size=15, color=WHITE)

# 下方合并箭头
add_down_arrow(slide, Inches(3.0), Inches(6.8), Inches(0.4), Inches(0.4), ACCENT)
add_down_arrow(slide, Inches(9.5), Inches(6.8), Inches(0.4), Inches(0.4), ACCENT)

add_text_box(slide, Inches(3.5), Inches(6.8), Inches(6), Inches(0.5),
             "两个方案并行输出，结果通过 State 合并，交给合规检查器", font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 第7页：合规检查与决策
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "合规检查 + 最终决策", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Pt(80), Pt(4), fill_color=ACCENT)

# 左：合规检查
card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.2), border_color=RED)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(4.8), Inches(4.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "合规检查器"
p.font.size = Pt(24)
p.font.color.rgb = RED
p.font.bold = True

comp_items = [
    "对照行业标准审查两个方案：",
    "  •  SY 5069-2017（常规井标准）",
    "  •  SY/T 5587.12-2018（打捞规程）",
    "  •  SYT 6987-2024（水平井标准）",
    "",
    "置信度计算规则：",
    "  •  基准分：0.72",
    "  •  缺失信息：每个 -0.03（上限 -0.25）",
    "  •  无相似案例：-0.12",
    "  •  最低不低于 0.35",
]
for c in comp_items:
    add_bullet_text(tf, f"  {c}", font_size=15, color=WHITE)

# 右：决策
card2 = add_shape(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(5.2), border_color=ACCENT2)
txBox2 = slide.shapes.add_textbox(Inches(7.5), Inches(2.0), Inches(5.0), Inches(4.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "决策 Maker"
p2.font.size = Pt(24)
p2.font.color.rgb = ACCENT2
p2.font.bold = True

dec_items = [
    "综合所有信息，生成最终方案：",
    "",
    "输入来源：",
    "  •  激进方案 + 保守方案",
    "  •  合规审查报告",
    "  •  证据链（案例 + 标准引用）",
    "",
    "输出结构：",
    "  •  事故概要 + 缺失信息",
    "  •  主路径（保守）+ 加速条件",
    "  •  快速恢复替代方案",
    "  •  判定节点 + 应急程序",
    "  •  风险提示 + 参考文献",
]
for d in dec_items:
    add_bullet_text(tf2, f"  {d}", font_size=15, color=WHITE)

# ════════════════════════════════════════════
# 第8页：多轮对话扩展
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "进阶：多轮对话图（Conversation Graph）", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Pt(80), Pt(4), fill_color=ACCENT)

# 路由器
make_agent_card(slide, Inches(5.5), Inches(1.6), Inches(2.5), Inches(1.0),
                "Router 路由器", "分类用户意图", ACCENT2, ACCENT2)

# 三个分支
branches = [
    ("solve 求解", Inches(0.5), "重新评估 / 修订方案\n触发完整辩论流程", ORANGE),
    ("explain 解释", Inches(4.8), "回答关于方案的追问\n不修改方案", ACCENT3),
    ("finalize 定稿", Inches(9.2), "锁定并导出方案\n结束对话", ACCENT),
]

for title, left, desc, color in branches:
    make_agent_card(slide, left, Inches(3.2), Inches(3.5), Inches(1.8), title, desc, color, color)

# solve 的子流程
add_text_box(slide, Inches(0.5), Inches(5.3), Inches(6), Inches(0.4),
             "solve 分支详细流程：", font_size=16, color=ORANGE, bold=True)
add_text_box(slide, Inches(0.5), Inches(5.8), Inches(10), Inches(1.2),
             "state_update → case_refresh → debate（模拟多轮辩论）→ compliance → revision（修订方案）→ archive",
             font_size=15, color=WHITE)

# ════════════════════════════════════════════
# 第9页：关键设计模式
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "关键设计模式总结", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Pt(80), Pt(4), fill_color=ACCENT)

patterns = [
    ("辩论模式\nDebate Pattern",
     "两个对立视角的 Agent 生成方案\n决策者综合双方观点\n模拟人类专家的讨论过程",
     ORANGE),
    ("证据链\nEvidence Chain",
     "每个 Agent 输出时追加 evidence\nAnnotated[List, add] 累加器\n最终方案引用完整证据链",
     ACCENT),
    ("LLM + 规则混合\nHybrid Approach",
     "每个 Agent 有确定性 fallback\nLLM 可选增强质量\n系统在无 LLM 时仍可运行",
     ACCENT3),
    ("Wiki 知识库\nRAG Pattern",
     "本地 Markdown 知识库\n加权关键词匹配检索\n标准文档 & 历史案例",
     ACCENT2),
]

for i, (title, desc, color) in enumerate(patterns):
    left = Inches(0.6 + i * 3.15)
    card = add_shape(slide, left, Inches(1.8), Inches(2.9), Inches(5.2), border_color=color)
    add_text_box(slide, left + Inches(0.2), Inches(2.0), Inches(2.5), Inches(1.0),
                 title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(3.3), Inches(2.5), Inches(3.2),
                 desc, font_size=14, color=WHITE)

# ════════════════════════════════════════════
# 第10页：代码结构一览
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "项目代码结构", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Pt(80), Pt(4), fill_color=ACCENT)

card = add_shape(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5.2), border_color=ACCENT)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.0), Inches(4.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "src/"
p.font.size = Pt(20)
p.font.color.rgb = ACCENT
p.font.bold = True

code_structure = [
    "├── agents/",
    "│   ├── accident_parser.py      # 事故信息提取节点",
    "│   ├── case_matcher.py         # 历史案例匹配节点",
    "│   ├── aggressive_plan.py      # 激进方案生成节点",
    "│   ├── conservative_plan.py    # 保守方案生成节点",
    "│   ├── compliance_checker.py   # 合规审查节点",
    "│   └── decision_maker.py       # 综合决策 + 归档节点",
    "├── graph.py                    # LangGraph 状态图定义（核心编排）",
    "├── state.py                    # AgentState / ConversationState 定义",
    "├── main.py                     # CLI 入口",
    "├── web_api.py                  # FastAPI Web 服务",
    "├── llm_client.py               # LLM 客户端（OpenAI / Anthropic / GLM）",
    "└── disposition_graph.py        # 知识图谱可视化",
]
for line in code_structure:
    add_bullet_text(tf, f"  {line}", font_size=14, color=WHITE)

# ════════════════════════════════════════════
# 第11页：总结
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(1),
             "总结与要点回顾", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(2.0), Inches(13.333), Pt(3), fill_color=ACCENT)

takeaways = [
    ("1", "多智能体 = 分工协作", "每个 Agent 有明确职责，通过共享 State 协同工作", ACCENT),
    ("2", "LangGraph 是编排引擎", "定义节点和边，支持串行、并行、条件路由", ACCENT2),
    ("3", "并行 + 辩论 = 更好决策", "两个对立方案并行生成，综合决策更全面", ORANGE),
    ("4", "State 是信息桥梁", "所有 Agent 通过 State 读写数据，Evidence 累加器实现证据链", ACCENT3),
    ("5", "LLM 是可选增强", "每个 Agent 都有确定性 fallback，系统鲁棒性强", RED),
]

for i, (num, title, desc, color) in enumerate(takeaways):
    top = Inches(2.4 + i * 0.95)
    # 编号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), top, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, Inches(1.5), top + Inches(0.08), Inches(0.6), Inches(0.5),
                 num, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.4), top + Inches(0.02), Inches(3), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, Inches(5.5), top + Inches(0.08), Inches(6.5), Inches(0.5),
                 desc, font_size=16, color=GRAY)

# 底部
add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
             "谢谢！欢迎提问交流", font_size=24, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = os.path.join(os.path.dirname(__file__), "多智能体协作原理.pptx")
prs.save(output_path)
print(f"PPT 已保存至: {output_path}")
